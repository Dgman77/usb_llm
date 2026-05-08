"""
rag.py — Adaptive Hybrid RAG (BM25 + FAISS) with multi-format + image support
Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, HTML, MD, JSON, XML, RTF, Images
Session-only storage — all data lives in RAM, cleared on server restart
"""

import os
import re
import io
import csv
import json
import base64
import html as html_mod
import xml.etree.ElementTree as ET
import numpy as np

from rank_bm25 import BM25Okapi
import faiss
import fitz                           # PDF + images → PyMuPDF
from docx import Document as DocxDoc  # DOCX → python-docx

# ── Optional imports (graceful fallback) ──────────────────
try:
    import openpyxl
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

try:
    from bs4 import BeautifulSoup
    _HAS_BS4 = True
except ImportError:
    _HAS_BS4 = False

try:
    from striprtf.striprtf import rtf_to_text
    _HAS_RTF = True
except ImportError:
    _HAS_RTF = False

try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False


# ── State (in-memory only) ────────────────────────────────
_chunks = []
_tokenized = []
_chunk_doc = []
_chunk_page = []
_chunk_type = []        # "text" | "image_meta"
_doc_names = []
_images = {}            # doc_name → [{data, ext, page, desc}, ...]
_vocab = {}

_bm25 = None
_faiss_index = None


# ── Tokenize ──────────────────────────────────────────────
def _tokenize(text: str):
    return re.findall(r"\b[a-z]{2,}\b", text.lower())


# ── Chunking with overlap ────────────────────────────────
def chunk_text(text: str, max_chars: int = 600, overlap: int = 80):
    paragraphs = re.split(r"\n\s*\n", text)
    chunks = []

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(para) <= max_chars:
            chunks.append(para)
        else:
            sentences = re.split(r"(?<=[.?!])\s+", para)
            current = ""
            for sent in sentences:
                sent = sent.strip()
                if not sent:
                    continue
                if len(current) + len(sent) + 1 <= max_chars:
                    current += (" " + sent) if current else sent
                else:
                    if current:
                        chunks.append(current.strip())
                    if overlap > 0 and len(current) > overlap:
                        current = current[-overlap:] + " " + sent
                    else:
                        current = sent
            if current:
                chunks.append(current.strip())

    return chunks if chunks else [text[:max_chars]]


# ── Build Hybrid Index ────────────────────────────────────
def _rebuild():
    global _bm25, _faiss_index, _tokenized, _vocab

    if not _chunks:
        _bm25 = None
        _faiss_index = None
        return

    _tokenized = [_tokenize(c) for c in _chunks]
    _bm25 = BM25Okapi(_tokenized)

    vocab = {}
    for tokens in _tokenized:
        for t in tokens:
            if t not in vocab:
                vocab[t] = len(vocab)

    dim = max(len(vocab), 1)
    mat = np.zeros((len(_chunks), dim), dtype="float32")

    for i, tokens in enumerate(_tokenized):
        for t in tokens:
            if t in vocab:
                mat[i][vocab[t]] += 1
        norm = np.linalg.norm(mat[i])
        if norm > 0:
            mat[i] /= norm

    _faiss_index = faiss.IndexFlatIP(dim)
    _faiss_index.add(mat)
    _vocab = vocab


# ═══════════════════════════════════════════════════════════
#  EXTRACTORS — each returns (pages, images)
#  pages  = [(page_num, text), ...]
#  images = [{data, ext, page, desc}, ...]   (may be empty)
# ═══════════════════════════════════════════════════════════

def _wrap_pages(pages):
    """Helper: return (pages, []) for extractors that have no images."""
    return pages, []


def extract_pdf(fb: bytes):
    doc = fitz.open(stream=fb, filetype="pdf")
    pages = [(i + 1, p.get_text()) for i, p in enumerate(doc)]
    images = []
    for i, page in enumerate(doc):
        for idx, img in enumerate(page.get_images(full=True)):
            try:
                bi = doc.extract_image(img[0])
                if bi:
                    images.append({
                        "data": bi["image"], "ext": bi["ext"], "page": i + 1,
                        "desc": f"Image {idx+1} page {i+1} ({bi.get('width',0)}x{bi.get('height',0)} {bi['ext']})"
                    })
            except Exception:
                pass
    return pages, images


def extract_docx(fb: bytes):
    doc = DocxDoc(io.BytesIO(fb))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    pages, buf, pn = [], "", 1
    for para in paragraphs:
        if len(buf) + len(para) > 1500:
            pages.append((pn, buf.strip())); buf = para + "\n"; pn += 1
        else:
            buf += para + "\n"
    if buf.strip():
        pages.append((pn, buf.strip()))
    images = []
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                ext = os.path.splitext(rel.target_part.partname)[1].lstrip(".")
                images.append({"data": rel.target_part.blob, "ext": ext, "page": 1,
                               "desc": f"Embedded image ({ext})"})
    except Exception:
        pass
    return pages, images


def extract_txt(fb: bytes):
    text = fb.decode("utf-8", errors="replace")
    lines = text.splitlines()
    pages, buf, pn = [], "", 1
    for line in lines:
        if len(buf) + len(line) > 1500:
            pages.append((pn, buf.strip())); buf = line + "\n"; pn += 1
        else:
            buf += line + "\n"
    if buf.strip():
        pages.append((pn, buf.strip()))
    return _wrap_pages(pages)


def extract_csv(fb: bytes):
    text = fb.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return _wrap_pages([])
    header = rows[0] if rows else []
    buf, pn, pages = "", 1, []
    for row in rows:
        line = " | ".join(row)
        if len(buf) + len(line) > 1500:
            pages.append((pn, buf.strip())); buf = line + "\n"; pn += 1
        else:
            buf += line + "\n"
    if buf.strip():
        pages.append((pn, buf.strip()))
    return _wrap_pages(pages)


def extract_xlsx(fb: bytes):
    if not _HAS_XLSX:
        raise ValueError("openpyxl not installed — cannot read .xlsx")
    wb = openpyxl.load_workbook(io.BytesIO(fb), read_only=True, data_only=True)
    pages = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        buf, pn = f"Sheet: {sheet}\n", 1
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(str(c) if c is not None else "" for c in row)
            if len(buf) + len(line) > 1500:
                pages.append((pn, buf.strip())); buf = line + "\n"; pn += 1
            else:
                buf += line + "\n"
        if buf.strip():
            pages.append((pn, buf.strip()))
    wb.close()
    return _wrap_pages(pages)


def extract_pptx(fb: bytes):
    if not _HAS_PPTX:
        raise ValueError("python-pptx not installed — cannot read .pptx")
    prs = Presentation(io.BytesIO(fb))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
        pages.append((i + 1, "\n".join(parts)))
    return _wrap_pages(pages)


def extract_html(fb: bytes):
    text = fb.decode("utf-8", errors="replace")
    if _HAS_BS4:
        soup = BeautifulSoup(text, "html.parser")
        clean = soup.get_text(separator="\n", strip=True)
    else:
        clean = re.sub(r"<[^>]+>", " ", text)
        clean = html_mod.unescape(clean)
    return extract_txt(clean.encode("utf-8"))


def extract_markdown(fb: bytes):
    return extract_txt(fb)


def extract_json(fb: bytes):
    text = fb.decode("utf-8", errors="replace")
    try:
        data = json.loads(text)
        pretty = json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        pretty = text
    return extract_txt(pretty.encode("utf-8"))


def extract_xml(fb: bytes):
    try:
        root = ET.fromstring(fb)
        parts = []
        for elem in root.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            txt = (elem.text or "").strip()
            if txt:
                parts.append(f"{tag}: {txt}")
        return extract_txt("\n".join(parts).encode("utf-8"))
    except ET.ParseError:
        return extract_txt(fb)


def extract_rtf(fb: bytes):
    if not _HAS_RTF:
        raise ValueError("striprtf not installed — cannot read .rtf")
    text = rtf_to_text(fb.decode("utf-8", errors="replace"))
    return extract_txt(text.encode("utf-8"))


def extract_image(fb: bytes, filename: str):
    """Extract metadata (and any text) from an image file."""
    desc_parts = [f"Image file: {filename}"]
    # Try PIL for dimensions
    if _HAS_PIL:
        try:
            img = Image.open(io.BytesIO(fb))
            w, h = img.size
            desc_parts.append(f"Dimensions: {w}x{h}")
            desc_parts.append(f"Mode: {img.mode}")
        except Exception:
            pass
    # Try PyMuPDF to extract any embedded text (OCR layer)
    ocr_text = ""
    try:
        ext = os.path.splitext(filename)[1].lstrip(".")
        doc = fitz.open(stream=fb, filetype=ext)
        if doc.page_count > 0:
            ocr_text = doc[0].get_text().strip()
    except Exception:
        pass

    meta_text = " | ".join(desc_parts)
    if ocr_text:
        meta_text += "\nExtracted text:\n" + ocr_text

    pages = [(1, meta_text)]
    images = [{"data": fb, "ext": os.path.splitext(filename)[1].lstrip("."),
               "page": 1, "desc": desc_parts[0]}]
    return pages, images


# ── Format router ─────────────────────────────────────────

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".svg"}

_EXTRACTOR_MAP = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".txt": extract_txt,
    ".csv": extract_csv,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".md": extract_markdown,
    ".json": extract_json,
    ".xml": extract_xml,
    ".rtf": extract_rtf,
    ".log": extract_txt,
    ".py": extract_txt,
    ".js": extract_txt,
    ".ts": extract_txt,
    ".css": extract_txt,
    ".yaml": extract_txt,
    ".yml": extract_txt,
    ".ini": extract_txt,
    ".cfg": extract_txt,
    ".toml": extract_txt,
}


def _get_extractor(filename: str):
    ext = os.path.splitext(filename.lower())[1]
    if ext in _IMAGE_EXTS:
        return lambda fb: extract_image(fb, filename)
    return _EXTRACTOR_MAP.get(ext)


def get_supported_formats() -> list[str]:
    fmts = sorted(set(list(_EXTRACTOR_MAP.keys()) + list(_IMAGE_EXTS)))
    return fmts


# ── Add Document ──────────────────────────────────────────
def add_document(file_bytes: bytes, filename: str) -> int:
    global _chunks, _chunk_doc, _chunk_page, _chunk_type, _doc_names

    if filename in _doc_names:
        return 0

    extractor = _get_extractor(filename)
    if extractor is None:
        supported = ", ".join(get_supported_formats())
        raise ValueError(f"Unsupported file type: {filename!r}. Supported: {supported}")

    result = extractor(file_bytes)
    # Extractors return either (pages, images) or just pages
    if isinstance(result, tuple) and len(result) == 2 and isinstance(result[1], list):
        pages, images = result
    else:
        pages, images = result, []

    count = 0
    for page_num, text in pages:
        if not text.strip():
            continue
        for chunk in chunk_text(text):
            _chunks.append(chunk)
            _chunk_doc.append(filename)
            _chunk_page.append(page_num)
            _chunk_type.append("text")
            count += 1

    # Store image metadata as searchable chunks
    if images:
        _images[filename] = images
        for img_info in images:
            meta_chunk = f"[Image in {filename}] {img_info['desc']}"
            _chunks.append(meta_chunk)
            _chunk_doc.append(filename)
            _chunk_page.append(img_info.get("page", 1))
            _chunk_type.append("image_meta")
            count += 1

    _doc_names.append(filename)
    _rebuild()
    return count


# ── Remove Document ───────────────────────────────────────
def remove_document(filename: str) -> bool:
    global _chunks, _chunk_doc, _chunk_page, _chunk_type, _doc_names

    if filename not in _doc_names:
        return False

    keep = [i for i in range(len(_chunks)) if _chunk_doc[i] != filename]
    _chunks = [_chunks[i] for i in keep]
    _chunk_doc = [_chunk_doc[i] for i in keep]
    _chunk_page = [_chunk_page[i] for i in keep]
    _chunk_type = [_chunk_type[i] for i in keep]
    _doc_names.remove(filename)
    _images.pop(filename, None)

    _rebuild()
    return True


# ── Clear All ─────────────────────────────────────────────
def clear_all():
    global _chunks, _tokenized, _chunk_doc, _chunk_page, _chunk_type
    global _doc_names, _bm25, _faiss_index, _images
    _chunks.clear(); _tokenized.clear(); _chunk_doc.clear()
    _chunk_page.clear(); _chunk_type.clear(); _doc_names.clear()
    _images.clear()
    _bm25 = None
    _faiss_index = None


# ── Query Vector ──────────────────────────────────────────
def _vectorize_query(tokens: list[str]) -> np.ndarray:
    vec = np.zeros(len(_vocab), dtype="float32")
    for t in tokens:
        if t in _vocab:
            vec[_vocab[t]] += 1
    norm = np.linalg.norm(vec)
    return vec / norm if norm > 0 else vec


# ── Hybrid Search (backward-compatible) ───────────────────
def search(query: str, top_k: int = 5) -> str:
    result = search_adaptive(query, top_k)
    return result["context"]


# ── Adaptive Search (new) ─────────────────────────────────
def search_adaptive(query: str, top_k: int = 5) -> dict:
    """
    Returns {
        context:    str   — retrieved text
        confidence: float — 0.0–1.0
        sources:    list  — [{doc, page, score}, ...]
        has_images: bool  — whether relevant images exist
        suggestion: str   — hint when confidence is low
    }
    """
    empty = {"context": "", "confidence": 0.0, "sources": [],
             "has_images": False, "suggestion": ""}

    if not _chunks:
        return empty

    tokens = _tokenize(query)
    if not tokens:
        return empty

    # BM25 scores
    bm25_scores = np.array(_bm25.get_scores(tokens))

    # FAISS scores
    q_vec = _vectorize_query(tokens)
    faiss_scores, faiss_idx = _faiss_index.search(
        q_vec.reshape(1, -1), min(20, len(_chunks))
    )

    # Combine
    combined = {}
    for score, idx in zip(faiss_scores[0], faiss_idx[0]):
        if idx >= 0:
            combined[idx] = combined.get(idx, 0) + float(score) * 0.5
    for i, s in enumerate(bm25_scores):
        combined[i] = combined.get(i, 0) + float(s) * 0.5

    ranked = sorted(combined.items(), key=lambda x: x[1], reverse=True)

    # Dedup & select
    selected, seen = [], set()
    for idx, score in ranked:
        key = _chunks[idx][:100]
        if key in seen:
            continue
        selected.append((idx, score))
        seen.add(key)
        if len(selected) >= top_k:
            break

    if not selected:
        return empty

    # ── Confidence scoring ────────────────────────────────
    top_score = selected[0][1]
    # Normalize: BM25 scores vary widely; use relative scoring
    bm25_max = float(bm25_scores.max()) if bm25_scores.max() > 0 else 1.0
    norm_top = top_score / (bm25_max + 0.5)  # rough normalization

    # Query coverage: how many query tokens appear in top results
    top_text = " ".join(_chunks[idx] for idx, _ in selected[:3]).lower()
    coverage = sum(1 for t in tokens if t in top_text) / max(len(tokens), 1)

    # Score gap: is top result clearly better?
    if len(selected) >= 2:
        gap = (selected[0][1] - selected[1][1]) / max(selected[0][1], 0.01)
    else:
        gap = 1.0

    confidence = min(1.0, (norm_top * 0.4 + coverage * 0.4 + gap * 0.2))

    # ── Build context ─────────────────────────────────────
    MAX_CHARS = 2500
    final = ""
    sources = []
    has_images = False

    for idx, score in selected:
        part = f"[Document: {_chunk_doc[idx]}, page {_chunk_page[idx]}]\n{_chunks[idx]}\n\n"
        if len(final) + len(part) > MAX_CHARS:
            break
        final += part
        sources.append({"doc": _chunk_doc[idx], "page": _chunk_page[idx],
                         "score": round(score, 4)})
        if _chunk_type[idx] == "image_meta":
            has_images = True

    # Check if any document has images
    if not has_images and _images:
        has_images = True

    # ── Suggestion for low confidence ─────────────────────
    suggestion = ""
    if confidence < 0.3 and final:
        doc_names = list(set(s["doc"] for s in sources))
        suggestion = (f"The query may not be directly covered, but related "
                      f"content was found in: {', '.join(doc_names)}")

    return {
        "context": final.strip(),
        "confidence": round(confidence, 3),
        "sources": sources,
        "has_images": has_images,
        "suggestion": suggestion,
    }


# ── Stats ─────────────────────────────────────────────────
def get_stats() -> dict:
    return {
        "total_chunks": len(_chunks),
        "documents": list(_doc_names),
        "doc_count": len(_doc_names),
        "image_count": sum(len(v) for v in _images.values()),
        "supported_formats": get_supported_formats(),
        "session_only": True,
    }


def get_all_content(max_chars: int = 3500) -> str:
    if not _chunks:
        return ""
    result = ""
    for i, chunk in enumerate(_chunks):
        part = f"[Document: {_chunk_doc[i]}, page {_chunk_page[i]}]\n{chunk}\n\n"
        if len(result) + len(part) > max_chars:
            remaining = max_chars - len(result)
            if remaining > 50:
                result += part[:remaining]
            break
        result += part
    return result.strip()


def has_documents() -> bool:
    return bool(_chunks)


def get_image_count() -> int:
    return sum(len(v) for v in _images.values())
